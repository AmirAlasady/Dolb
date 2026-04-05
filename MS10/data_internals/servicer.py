# MS10/data_internals/servicer.py

import grpc
import logging
import fitz  # PyMuPDF
import docx  # python-docx
import pptx  # python-pptx
from google.protobuf.struct_pb2 import Struct
from django.core.files.storage import default_storage
from rest_framework.exceptions import PermissionDenied

from . import data_pb2, data_pb2_grpc
from data.models import StoredFile

logging.basicConfig(level=logging.INFO, format='%(asctime)s - MS10-gRPC - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


class DataServicer(data_pb2_grpc.DataServiceServicer):
    """
    Implements the full DataService gRPC interface, including content parsing
    for various document types and metadata retrieval.
    """

    def GetFileMetadata(self, request, context):
        """
        Handles the fast metadata request from MS5 for validation.
        """
        logger.info(f"gRPC [GetFileMetadata]: Received request for {len(request.file_ids)} files from user {request.user_id}.")
        try:
            # Query for all requested files owned by the user in a single database hit.
            files = StoredFile.objects.filter(id__in=request.file_ids, owner_id=request.user_id)
            
            # This check handles both "not found" and "permission denied" implicitly.
            if files.count() != len(request.file_ids):
                context.set_code(grpc.StatusCode.NOT_FOUND)
                context.set_details("One or more files not found or you do not have permission to access them.")
                logger.warning(f"gRPC [GetFileMetadata]: Validation failed for user {request.user_id}.")
                return data_pb2.GetFileMetadataResponse()

            metadata_list = [
                data_pb2.FileMetadata(file_id=str(f.id), mimetype=f.mimetype)
                for f in files
            ]
            logger.info(f"gRPC [GetFileMetadata]: Successfully validated and found metadata for {len(metadata_list)} files.")
            return data_pb2.GetFileMetadataResponse(metadata=metadata_list)

        except Exception as e:
            logger.error(f"gRPC [GetFileMetadata]: Internal error: {e}", exc_info=True)
            context.set_code(grpc.StatusCode.INTERNAL)
            context.set_details(f"An internal error occurred: {e}")
            return data_pb2.GetFileMetadataResponse()

    def GetFileContent(self, request, context):
        """
        Handles the content retrieval request from MS6.
        """
        logger.info(f"gRPC [GetFileContent]: Received request for file {request.file_id} from user {request.user_id}.")
        try:
            file_instance = StoredFile.objects.get(id=request.file_id)
            if str(file_instance.owner_id) != str(request.user_id):
                raise PermissionDenied()

            raw_file_stream = default_storage.open(file_instance.storage_path)
            content_payload = self._parse_content(raw_file_stream, file_instance.mimetype, file_instance.storage_path)
            
            proto_struct = Struct()
            proto_struct.update(content_payload)
            
            logger.info(f"gRPC [GetFileContent]: Successfully parsed and returning content for file {request.file_id}.")
            return data_pb2.GetFileContentResponse(file_id=str(file_instance.id), content=proto_struct)

        except StoredFile.DoesNotExist:
            context.set_code(grpc.StatusCode.NOT_FOUND)
            context.set_details("File not found.")
            return data_pb2.GetFileContentResponse()
        except PermissionDenied:
            context.set_code(grpc.StatusCode.PERMISSION_DENIED)
            context.set_details("Permission denied to access this file.")
            return data_pb2.GetFileContentResponse()
        except Exception as e:
            logger.error(f"gRPC [GetFileContent]: Internal error: {e}", exc_info=True)
            context.set_code(grpc.StatusCode.INTERNAL)
            context.set_details(f"An internal error occurred during content retrieval: {e}")
            return data_pb2.GetFileContentResponse()

    # --- THE NEW PARSING ROUTER AND HELPERS ---

    def _parse_content(self, file_stream, mimetype, storage_path):
        """Routes the file stream to the correct parser based on its mimetype."""
        logger.info(f"Parsing content for mimetype: {mimetype}")
        
        # DOCX
        if mimetype == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            return self._parse_docx(file_stream)
        
        # PPTX
        elif mimetype == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            return self._parse_pptx(file_stream)
            
        # PDF
        elif mimetype == 'application/pdf':
            return self._parse_pdf(file_stream)
        
        # Plain Text
        elif mimetype.startswith('text/'):
            return {"type": "text_content", "content": file_stream.read().decode('utf-8')}
        
        # Image
        elif mimetype.startswith('image/'):
            url = default_storage.url(storage_path)
            return {"type": "image_url", "url": url}
            
        # Fallback for unsupported types
        else:
            logger.warning(f"Unsupported mimetype '{mimetype}' encountered. Cannot parse content.")
            return {"type": "unsupported", "content": f"File type '{mimetype}' is not supported for content parsing."}

    def _parse_docx(self, file_stream):
        """Extracts text from a .docx file stream."""
        try:
            document = docx.Document(file_stream)
            full_text = [para.text for para in document.paragraphs]
            return {"type": "text_content", "content": '\n'.join(full_text)}
        except Exception as e:
            logger.error(f"Failed to parse DOCX file: {e}", exc_info=True)
            return {"type": "error", "content": "Failed to parse the DOCX file."}

    def _parse_pptx(self, file_stream):
        """Extracts text from shapes in a .pptx file stream."""
        try:
            presentation = pptx.Presentation(file_stream)
            full_text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
            return {"type": "text_content", "content": '\n\n'.join(full_text)}
        except Exception as e:
            logger.error(f"Failed to parse PPTX file: {e}", exc_info=True)
            return {"type": "error", "content": "Failed to parse the PPTX file."}
        
    def _parse_pdf(self, file_stream):
        """Extracts text from a .pdf file stream."""
        try:
            text = ""
            with fitz.open(stream=file_stream.read(), filetype="pdf") as doc:
                for page in doc:
                    text += page.get_text()
            return {"type": "text_content", "content": text}
        except Exception as e:
            logger.error(f"Failed to parse PDF file: {e}", exc_info=True)
            return {"type": "error", "content": "Failed to parse the PDF file."}